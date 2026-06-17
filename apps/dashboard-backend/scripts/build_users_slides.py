"""Build 발표_37-40_users.pptx — 사용자 관리(RBAC) 시연 4장.

시연 영상이 3개(생성+첫 로그인 / 권한 수정 / 삭제)로 나뉘므로
기존 발표_37-38_users.pptx 의 단일 영상 슬라이드를 영상 3장으로 분리한다.

스타일 기준: 발표_37-38_users.pptx(초록 테마) + 발표_31-34_fire.pptx 영상 슬라이드 geometry.
내용 기준: docs/demo/03_user_account_scenario.md, docs/report/05_발표_틀.md 37~38.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (user mgmt = green) ----
NAVY = RGBColor(0x1F, 0x38, 0x64)
BODY = RGBColor(0x40, 0x40, 0x40)
SUB = RGBColor(0x59, 0x59, 0x59)
GREEN = RGBColor(0x54, 0x82, 0x35)       # 강조색 초록 (원/화살표)
GREEN_DK = RGBColor(0x37, 0x56, 0x23)    # 라벨 진초록
GREEN_BG = RGBColor(0xEC, 0xF3, 0xE3)    # 연초록 카드 배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xBF, 0xBF, 0xBF)
CARD_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width = Emu(12191695)
prs.slide_height = Emu(6858000)
blank = prs.slide_layouts[6]


def _set_font(run, size, bold, color):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def _prep_tf(tf, anchor):
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    _prep_tf(tb.text_frame, anchor)
    for i, (text, size, bold, color) in enumerate(lines):
        p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        _set_font(r, size, bold, color)
    return tb


def add_rich(slide, x, y, w, h, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    _prep_tf(tb.text_frame, anchor)
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        _set_font(r, size, bold, color)
    return tb


def add_shape(slide, shape_type, x, y, w, h, fill, line=None, line_w=None):
    sp = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        if line_w:
            sp.line.width = Emu(line_w)
    sp.shadow.inherit = False
    return sp


def add_circle_num(slide, x, y, d, num, fill, size=26):
    sp = add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill)
    _prep_tf(sp.text_frame, MSO_ANCHOR.MIDDLE)
    p = sp.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    _set_font(r, size, True, WHITE)
    return sp


# ============================================================
# Slide 37 — 사용자 관리 시나리오 설명 (기존 슬라이드 재현)
# ============================================================
s = prs.slides.add_slide(blank)

add_text(s, 457200, 301752, 11274552, 758952,
         [("시연 ③ 사용자 관리 (권한 기반 접근 통제)", 40, True, NAVY)],
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 457200, 1170432, 11274552, 822960,
         [("로그인·세션은 Cognito, 공장 접근 권한은 RDS가 담당  →  백엔드가 둘을 매칭해 '권한에 맞는 공장만' 노출",
           16, False, SUB)],
         anchor=MSO_ANCHOR.MIDDLE)

cards = [
    ("1", "생성 + 첫 로그인", "( factory-a 만 노출 )",
     ["공장 관리자 계정 생성 →", "초대 메일·임시 비밀번호로 첫 로그인"]),
    ("2", "권한 수정", "( factory-a + factory-b )",
     ["공장 접근 권한 변경 →", "화면의 공장 범위가 즉시 변화"]),
    ("3", "삭제", "( 로그인 차단 )",
     ["Cognito 계정 + RDS 권한을", "함께 제거 → 재로그인 불가"]),
]
card_x = [457200, 4498848, 8531352]
for cx, (num, title, label, body) in zip(card_x, cards):
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, 2423160, 3200400, 2331720,
              GREEN_BG, line=GREEN, line_w=19050)
    add_circle_num(s, cx + 1243584, 2624328, 713232, num, GREEN)
    add_text(s, cx + 182880, 3447288, 2834640, 411480, [(title, 19, True, NAVY)])
    add_text(s, cx + 182880, 3867912, 2834640, 310896, [(label, 13, True, GREEN_DK)])
    add_text(s, cx + 228600, 4233672, 2743200, 457200,
             [(b, 12, False, BODY) for b in body])

for ax in (3758184, 7790688):
    add_shape(s, MSO_SHAPE.RIGHT_ARROW, ax, 3355848, 649224, 475488, GREEN)

bar = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 457200, 5522976, 11274552, 822960, NAVY)
_prep_tf(bar.text_frame, MSO_ANCHOR.MIDDLE)
p = bar.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "본사 관리자가 화면에서 공장별 접근 권한을 직접 운영한다"
_set_font(r, 18, True, WHITE)

# ============================================================
# Slide 38~40 — 영상 3장 (시연 영상이 3개로 분리됨)
# ============================================================
videos = [
    ("1", "생성 + 첫 로그인 — 초대 메일·비밀번호 변경",
     "생성 + 첫 로그인 : ",
     "공장 관리자(factory-a) 생성 → 초대 메일·임시 비밀번호 → 첫 로그인 시 비밀번호 변경 → factory-a 만 노출"),
    ("2", "권한 수정 — 공장 접근 범위 변경",
     "권한 수정 : ",
     "factory-a → factory-a + factory-b 변경·저장 → 신규 사용자 화면 새로고침 → 보이는 공장 범위 즉시 변경"),
    ("3", "삭제 — 계정 제거·로그인 차단",
     "삭제 : ",
     "계정 삭제 → 목록에서 사라짐 → 재로그인 불가 (Cognito 계정 + RDS 권한 함께 제거)"),
]
for num, title, cap_label, cap_body in videos:
    s = prs.slides.add_slide(blank)
    add_circle_num(s, 502920, 393192, 566928, num, GREEN, size=22)
    add_text(s, 1234440, 365760, 10469880, 621792,
             [(title, 26, True, NAVY)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # 영상 placeholder (16:9)
    add_shape(s, MSO_SHAPE.RECTANGLE, 2195931, 1325880, 7799831, 4389120,
              CARD_GRAY, line=CARD_BORDER, line_w=15875)
    add_shape(s, MSO_SHAPE.OVAL, 5615787, 2857499, 960120, 960120, GREEN)
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                             Emu(5912967), Emu(3108959), Emu(402336), Emu(457200))
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = WHITE
    tri.line.fill.background()
    tri.shadow.inherit = False
    add_text(s, 2195931, 4023359, 7799831, 365760,
             [("영상 자리 — PowerPoint에서 영상 삽입", 13, False, SUB)])
    add_rich(s, 457200, 5989320, 11274552, 566928,
             [(cap_label, 14, True, GREEN_DK), (cap_body, 14, False, BODY)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

out = "docs/report/발표_37-40_users.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))

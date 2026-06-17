"""Build 발표_35-36_cloud.pptx — 클라우드 에러(ECS 장애) 시연 슬라이드 2장.

스타일 기준: docs/report/발표_31-34_fire.pptx (동일 geometry/폰트), 단 파랑 테마.
내용 기준: docs/demo/02_ecs_test.md, docs/report/05_발표_틀.md 35~36.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (cloud-error = blue) ----
NAVY = RGBColor(0x1F, 0x38, 0x64)      # 제목 진한 네이비
BODY = RGBColor(0x40, 0x40, 0x40)      # 본문
SUB = RGBColor(0x59, 0x59, 0x59)       # 부제
BLUE = RGBColor(0x2E, 0x75, 0xB6)      # 강조색 파랑
BLUE_BG = RGBColor(0xE8, 0xF1, 0xFA)   # 연파랑 배경
POINT = RGBColor(0x1F, 0x4E, 0x79)     # 포인트 네이비
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xBF, 0xBF, 0xBF)
CARD_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width = Emu(12191695)
prs.slide_height = Emu(6858000)
blank = prs.slide_layouts[6]


def _set_font(run, size, bold, color):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    # 한글 폰트(eastasian) 강제 지정
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color) — 한 줄당 한 문단."""
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        _set_font(r, size, bold, color)
    return tb


def add_rich(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    """paras: list of list of (text, size, bold, color) — 문단 내 여러 run."""
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        for text, size, bold, color in runs:
            r = p.add_run()
            r.text = text
            _set_font(r, size, bold, color)
    return tb


def add_shape(slide, shape_type, x, y, w, h, fill, line=None, line_w=None):
    sp = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        if line_w:
            sp.line.width = Emu(line_w)
    sp.shadow.inherit = False
    return sp


def add_circle_num(slide, x, y, d, num, fill, size=26):
    sp = add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    _set_font(r, size, True, WHITE)
    return sp


# ============================================================
# Slide 35 — 클라우드 에러 시나리오 소개
# ============================================================
s = prs.slides.add_slide(blank)

add_text(s, 457200, 301752, 11274552, 758952,
         [("시연 ② 클라우드 에러 (ECS 장애)", 40, True, NAVY)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, 457200, 1170432, 11274552, 731520,
         [("관제 시스템 자체(대시보드 백엔드 ECS)에 장애가 나면?  →  Cloud Infra가 스스로 감지  →  Slack 알림  →  자동 복구",
           16, False, SUB)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

cards = [
    ("1", "장애 발생", "( 백엔드 ECS )",
     ["ECS 2대 중 1대를 일부러 중지", "목표 2 · 실행 1 (warning)"]),
    ("2", "자동 감지·알림", "( Cloud Infra )",
     ["화면은 살아 있는 채 이상만 표시", "Slack으로 이상 알림 전송"]),
    ("3", "자동 복구", "( ECS Auto-Recovery )",
     ["새 task 자동 기동", "정상(2/2)으로 복귀"]),
]
card_x = [457200, 4494276, 8531352]
for cx, (num, title, label, body) in zip(card_x, cards):
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, 2423160, 3200400, 2331720,
              BLUE_BG, line=BLUE, line_w=19050)
    add_circle_num(s, cx + 1243584, 2624328, 713232, num, BLUE)
    add_text(s, cx + 182880, 3447288, 2834640, 411480,
             [(title, 19, True, NAVY)])
    add_text(s, cx + 182880, 3867912, 2834640, 310896,
             [(label, 13, True, BLUE)])
    add_text(s, cx + 228600, 4233672, 2743200, 457200,
             [(b, 12, False, BODY) for b in body])

# 카드 사이 화살표
for ax in (3753612, 7790688):
    add_shape(s, MSO_SHAPE.RIGHT_ARROW, ax, 3351276, 649224, 475488, BLUE)

# 하단 네이비 메시지 바
bar = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 457200, 5522976, 11274552, 822960, NAVY)
tf = bar.text_frame
tf.word_wrap = True
tf.margin_left = 0
tf.margin_right = 0
tf.margin_top = 0
tf.margin_bottom = 0
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "현장 공장뿐 아니라 관제 시스템의 건강 상태도 모니터링하고, 문제가 나면 자동으로 복구합니다"
_set_font(r, 18, True, WHITE)

# ============================================================
# Slide 36 — 클라우드 에러 영상
# ============================================================
s = prs.slides.add_slide(blank)

add_circle_num(s, 502920, 393192, 566928, "2", BLUE, size=22)
add_text(s, 1234440, 365760, 10424160, 621792,
         [("ECS 장애 → 자동 복구", 26, True, NAVY)],
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# 영상 placeholder (16:9 비율)
add_shape(s, MSO_SHAPE.RECTANGLE, 2195931, 1325880, 7799831, 4389120,
          CARD_GRAY, line=CARD_BORDER, line_w=15875)
add_shape(s, MSO_SHAPE.OVAL, 5615787, 2857499, 960120, 960120, BLUE)
tri = slide_tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                     Emu(5912967), Emu(3108959), Emu(402336), Emu(457200))
tri.rotation = 90
tri.fill.solid()
tri.fill.fore_color.rgb = WHITE
tri.line.fill.background()
tri.shadow.inherit = False
add_text(s, 2195931, 4023359, 7799831, 365760,
         [("영상 자리 — PowerPoint에서 영상 삽입", 13, False, SUB)])

# 하단 캡션 2줄 (요지만)
add_rich(s, 457200, 5780000, 11274552, 980000,
         [
             [("장애   ", 13, True, BLUE),
              ("ECS 1대 강제 종료 → 콘솔 2/2→1/2 → Dashboard '주의' → ALB target 하락 → Slack 알림",
               13, False, BODY)],
             [("복구   ", 13, True, BLUE),
              ("ECS 자동 복구 → Dashboard '정상' → (옛 task DRAINING 동안 ALB 잠깐 '주의') → ALB도 정상 target 2개 복귀",
               13, False, BODY)],
         ],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

out = "docs/report/발표_35-36_cloud.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
